"""
Generate a PowerPoint presentation explaining the PanNET bipartite GNN pipeline.
Visual-heavy, minimal text — diagrams, arrows, architecture figures.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors — white background theme
DARK_BG = RGBColor(0xFF, 0xFF, 0xFF)       # White background
BLACK = RGBColor(0x1A, 0x1A, 0x1A)         # Headlines and headers — always black
WHITE = RGBColor(0x1A, 0x1A, 0x1A)         # Body text (dark on white)
LIGHT_GRAY = RGBColor(0x55, 0x55, 0x66)    # Secondary/descriptive text
CARD_TEXT = RGBColor(0xFF, 0xFF, 0xFF)      # White text on colored cards
ACCENT_GREEN = RGBColor(0x27, 0xAE, 0x60)  # PanNET / tumor identification
ACCENT_RED = RGBColor(0xC0, 0x39, 0x2B)    # Stroma identification
ACCENT_BLUE = RGBColor(0x29, 0x80, 0xB9)   # Normal tissue identification
ACCENT_ORANGE = RGBColor(0xD3, 0x54, 0x00) # Warnings / key findings
ACCENT_PURPLE = RGBColor(0x7D, 0x3C, 0x98) # Model / architecture
DARK_CARD = RGBColor(0xF0, 0xF0, 0xF5)     # Light card background
MID_GRAY = RGBColor(0x7F, 0x8C, 0x8D)      # Muted elements / arrows
BORDER_GRAY = RGBColor(0xBB, 0xBB, 0xCC)   # Neutral card borders

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def add_dark_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_card(slide, left, top, width, height, fill_color=DARK_CARD, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    return shape


def add_arrow(slide, left, top, width, height, color=ACCENT_ORANGE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_down_arrow(slide, left, top, width, height, color=ACCENT_ORANGE):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_circle(slide, left, top, size, fill_color, text="", font_size=10):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = CARD_TEXT
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(0)
        tf.paragraphs[0].space_after = Pt(0)
    return shape


def add_multiline_text(slide, left, top, width, height, lines, font_size=14,
                       color=WHITE, bold=False, line_spacing=1.2, bullet=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = ("• " if bullet else "") + line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Calibri"
        p.space_after = Pt(font_size * (line_spacing - 1) + 2)
    return txBox


# ============================================================================
# SLIDE 1: Title
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_dark_bg(slide)

add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
             "Context-Aware Graph Neural Networks", font_size=40, bold=True, color=BLACK,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(2.9), Inches(10), Inches(0.8),
             "for PanNET Infiltration Pattern Scoring", font_size=28, color=BLACK,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(4.2), Inches(10), Inches(0.5),
             "End-to-End Pipeline: WSI → Bipartite GNN → Patient-Level IPS Prediction",
             font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Pipeline overview bar
stage_names = ["Stage 1\nWSI → Patches", "Stage 2\nAutoEncoder", "Stage 3\nGraph Build", "Stage 4\nCache (.pkl)", "Stage 5\nGIN Training"]
stage_colors = [ACCENT_BLUE, ACCENT_PURPLE, ACCENT_GREEN, ACCENT_ORANGE, ACCENT_RED]
bar_y = Inches(5.3)
bar_w = Inches(2.0)
bar_h = Inches(0.9)
bar_start = Inches(1.2)
gap = Inches(0.35)

for i, (name, col) in enumerate(zip(stage_names, stage_colors)):
    x = bar_start + i * (bar_w + gap)
    card = add_card(slide, x, bar_y, bar_w, bar_h, fill_color=col)
    tf = card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(12)
    p.font.color.rgb = CARD_TEXT
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(6)

    if i < 4:
        add_arrow(slide, x + bar_w, bar_y + Inches(0.3), Inches(0.35), Inches(0.3), color=MID_GRAY)


# ============================================================================
# SLIDE 2: Medical Background — What is PanNET IPS?
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
             "Medical Background: PanNET Infiltration Pattern Scoring",
             font_size=28, bold=True, color=BLACK)
add_text_box(slide, Inches(0.5), Inches(0.9), Inches(12), Inches(0.4),
             "How pathologists grade tumor infiltration at the tumor–NNP interface",
             font_size=14, color=LIGHT_GRAY)

# Grade table
grades = [
    ("Grade 1", "Fully demarcated", "Round/capsulated tumor, clean boundary"),
    ("Grade 2", "Mildly irregular", "Early capsular penetration, clusters attached"),
    ("Grade 3", "Satellite/projections", "Large nodules, connected to main mass"),
    ("Grade 4", "Peri-tumoral", "Small clusters nearby, no distant invasion"),
    ("Grade 5", "Non-demarcated", "Prominent infiltration, clusters far from tumor"),
]

table_top = Inches(1.5)
row_h = Inches(0.7)
for i, (grade, pattern, desc) in enumerate(grades):
    y = table_top + i * row_h
    # Grade badge
    badge_col = ACCENT_GREEN if i < 2 else (ACCENT_ORANGE if i < 4 else ACCENT_RED)
    card = add_card(slide, Inches(0.5), y, Inches(1.2), Inches(0.55), fill_color=badge_col)
    tf = card.text_frame
    p = tf.paragraphs[0]
    p.text = grade
    p.font.size = Pt(14)
    p.font.color.rgb = CARD_TEXT
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_text_box(slide, Inches(1.9), y, Inches(2.5), Inches(0.55), pattern,
                 font_size=14, bold=True, color=WHITE)
    add_text_box(slide, Inches(4.5), y, Inches(4.5), Inches(0.55), desc,
                 font_size=12, color=LIGHT_GRAY)

# IPS scoring on the right
ips_x = Inches(9.3)
add_text_box(slide, ips_x, Inches(1.5), Inches(3.5), Inches(0.5),
             "Patient-Level IPS", font_size=20, bold=True, color=BLACK)

add_text_box(slide, ips_x, Inches(2.1), Inches(3.8), Inches(0.4),
             "Sum grades from 3 representative WSIs:", font_size=12, color=LIGHT_GRAY)

ips_data = [
    ("IPS-A", "Sum ∈ [3, 6]", "Non/minimally infiltrative", ACCENT_GREEN),
    ("IPS-B", "Sum ∈ [7, 9]", "Moderately infiltrative", ACCENT_ORANGE),
    ("IPS-C", "Sum ∈ [10, 15]", "Highly infiltrative", ACCENT_RED),
]

for i, (name, range_str, desc, col) in enumerate(ips_data):
    y = Inches(2.7) + i * Inches(0.85)
    card = add_card(slide, ips_x, y, Inches(3.5), Inches(0.7), border_color=BORDER_GRAY)
    tf = card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"{name}  |  {range_str}"
    p.font.size = Pt(14)
    p.font.color.rgb = BLACK
    p.font.bold = True
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(11)
    p2.font.color.rgb = LIGHT_GRAY

# Key insight box
insight_y = Inches(5.5)
card = add_card(slide, Inches(0.5), insight_y, Inches(12.3), Inches(1.2), border_color=BORDER_GRAY)
tf = card.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "KEY INSIGHT"
p.font.size = Pt(14)
p.font.color.rgb = ACCENT_ORANGE
p.font.bold = True
p2 = tf.add_paragraph()
p2.text = "Infiltration is a SPATIAL property — it depends on HOW tumor interacts with surrounding tissue (NNP), not just what tumor cells look like individually. This is why GNNs are needed: they model spatial relationships between tissue patches."
p2.font.size = Pt(13)
p2.font.color.rgb = WHITE


# ============================================================================
# SLIDE 3: Pipeline Overview (full flow diagram)
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
             "End-to-End Pipeline Overview", font_size=28, bold=True, color=BLACK)

# Row 1: WSI → Patches → Features
row1_y = Inches(1.3)
steps_r1 = [
    ("WSI\n(.tiff)", "100K×100K px\nH&E stained", ACCENT_BLUE),
    ("Tissue\nSegmentation", "GrandQC\n≥60% tissue", MID_GRAY),
    ("Patch\nExtraction", "1024×1024 px\n@ 40x mag", ACCENT_BLUE),
    ("VirChow2\nEncoder", "Foundation model\n→ 1280-d vector", ACCENT_PURPLE),
    ("H5 File\nOutput", "features (N,1280)\ncoords (N,2)", ACCENT_GREEN),
]

step_w = Inches(2.1)
step_h = Inches(1.3)
start_x = Inches(0.4)
for i, (title, desc, col) in enumerate(steps_r1):
    x = start_x + i * (step_w + Inches(0.3))
    card = add_card(slide, x, row1_y, step_w, step_h, border_color=BORDER_GRAY)
    tf = card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(11)
    p2.font.color.rgb = LIGHT_GRAY
    p2.alignment = PP_ALIGN.CENTER

    if i < 4:
        add_arrow(slide, x + step_w, row1_y + Inches(0.45), Inches(0.3), Inches(0.25), MID_GRAY)

# Stage label
add_text_box(slide, Inches(0.4), row1_y - Inches(0.3), Inches(2), Inches(0.3),
             "STAGE 1", font_size=12, bold=True, color=ACCENT_BLUE)

# Down arrow
add_down_arrow(slide, Inches(6.3), row1_y + step_h, Inches(0.4), Inches(0.4), MID_GRAY)

# Row 2: AutoEncoder + Graph Construction
row2_y = Inches(3.3)
add_text_box(slide, Inches(0.4), row2_y - Inches(0.3), Inches(2), Inches(0.3),
             "STAGE 2 + 3", font_size=12, bold=True, color=ACCENT_PURPLE)

steps_r2 = [
    ("AutoEncoder\n(frozen)", "1280 → 768 → 512\n→ 256-d latent", ACCENT_PURPLE),
    ("Patch\nClassifier", "3-class MLP\nStroma/PanNET/Normal", ACCENT_ORANGE),
    ("Border\nDetection", "8-connected check\n+ NNP neighbors", ACCENT_GREEN),
    ("Bipartite\nGraph", "PanNET ↔ NNP edges\nChebyshev dist ≤ r", ACCENT_GREEN),
    ("Cache\nas .pkl", "PyG Data objects\nnode + edge attrs", MID_GRAY),
]

for i, (title, desc, col) in enumerate(steps_r2):
    x = start_x + i * (step_w + Inches(0.3))
    card = add_card(slide, x, row2_y, step_w, step_h, border_color=BORDER_GRAY)
    tf = card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(11)
    p2.font.color.rgb = LIGHT_GRAY
    p2.alignment = PP_ALIGN.CENTER

    if i < 4:
        add_arrow(slide, x + step_w, row2_y + Inches(0.45), Inches(0.3), Inches(0.25), MID_GRAY)

# Down arrow
add_down_arrow(slide, Inches(6.3), row2_y + step_h, Inches(0.4), Inches(0.4), MID_GRAY)

# Row 3: GIN Training
row3_y = Inches(5.3)
add_text_box(slide, Inches(0.4), row3_y - Inches(0.3), Inches(2), Inches(0.3),
             "STAGE 5", font_size=12, bold=True, color=ACCENT_RED)

steps_r3 = [
    ("GIN\nMessage Passing", "Sum aggregation\n+ residual connections", ACCENT_RED),
    ("PanNET-Only\nPooling", "global_add_pool\ntumor nodes only", ACCENT_RED),
    ("Regression\nHead", "MLP → single grade\nHuber loss (δ=2)", ACCENT_ORANGE),
    ("Patient\nAggregation", "3 slides → sum\n→ IPS (A/B/C)", ACCENT_ORANGE),
    ("Evaluation\nMetrics", "F1, QWK, MAE\n4-fold × 25 seeds", ACCENT_GREEN),
]

for i, (title, desc, col) in enumerate(steps_r3):
    x = start_x + i * (step_w + Inches(0.3))
    card = add_card(slide, x, row3_y, step_w, step_h, border_color=BORDER_GRAY)
    tf = card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(11)
    p2.font.color.rgb = LIGHT_GRAY
    p2.alignment = PP_ALIGN.CENTER

    if i < 4:
        add_arrow(slide, x + step_w, row3_y + Inches(0.45), Inches(0.3), Inches(0.25), MID_GRAY)


# ============================================================================
# SLIDE 4: Stage 1 — WSI to Patches in Detail
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(5), Inches(0.6),
             "Stage 1: WSI → Feature Extraction", font_size=28, bold=True, color=BLACK)

# WSI visualization (large rectangle)
wsi_card = add_card(slide, Inches(0.5), Inches(1.3), Inches(3.5), Inches(2.5), border_color=BORDER_GRAY)
tf = wsi_card.text_frame
p = tf.paragraphs[0]
p.text = "Whole Slide Image"
p.font.size = Pt(14)
p.font.color.rgb = ACCENT_BLUE
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

# Grid of patches inside WSI
for row in range(4):
    for col in range(5):
        px = Inches(0.7) + col * Inches(0.58)
        py = Inches(2.0) + row * Inches(0.4)
        colors = [ACCENT_GREEN, ACCENT_RED, ACCENT_BLUE, DARK_CARD]
        c = colors[(row + col) % 4] if (row + col) % 3 != 0 else DARK_BG
        add_card(slide, px, py, Inches(0.5), Inches(0.32), fill_color=c)

# Legend
add_circle(slide, Inches(0.5), Inches(4.1), Inches(0.25), ACCENT_GREEN)
add_text_box(slide, Inches(0.8), Inches(4.1), Inches(1), Inches(0.3), "PanNET", font_size=10, color=ACCENT_GREEN)
add_circle(slide, Inches(1.7), Inches(4.1), Inches(0.25), ACCENT_RED)
add_text_box(slide, Inches(2.0), Inches(4.1), Inches(1), Inches(0.3), "Stroma", font_size=10, color=ACCENT_RED)
add_circle(slide, Inches(2.8), Inches(4.1), Inches(0.25), ACCENT_BLUE)
add_text_box(slide, Inches(3.1), Inches(4.1), Inches(1), Inches(0.3), "Normal", font_size=10, color=ACCENT_BLUE)

# Arrow
add_arrow(slide, Inches(4.2), Inches(2.3), Inches(0.7), Inches(0.4), ACCENT_ORANGE)

# VirChow2 box
ae_card = add_card(slide, Inches(5.1), Inches(1.3), Inches(3.2), Inches(3.0), border_color=BORDER_GRAY)
tf = ae_card.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "VirChow2"
p.font.size = Pt(18)
p.font.color.rgb = ACCENT_PURPLE
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

lines = [
    ("Foundation Model", WHITE),
    ("Pre-trained on millions\nof pathology images", LIGHT_GRAY),
    ("", WHITE),
    ("Input: 1024×1024×3 patch", LIGHT_GRAY),
    ("Output: 1280-d vector", ACCENT_GREEN),
    ("", WHITE),
    ("Captures: cell morphology,", LIGHT_GRAY),
    ("tissue structure, staining", LIGHT_GRAY),
]
for text, col in lines:
    p2 = tf.add_paragraph()
    p2.text = text
    p2.font.size = Pt(11)
    p2.font.color.rgb = col
    p2.alignment = PP_ALIGN.CENTER

# Arrow
add_arrow(slide, Inches(8.5), Inches(2.3), Inches(0.7), Inches(0.4), ACCENT_ORANGE)

# H5 output
h5_card = add_card(slide, Inches(9.4), Inches(1.3), Inches(3.4), Inches(3.0), border_color=BORDER_GRAY)
tf = h5_card.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "H5 File (per slide)"
p.font.size = Pt(16)
p.font.color.rgb = ACCENT_GREEN
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

h5_lines = [
    ("features     (N, 1280)", ACCENT_GREEN),
    ("coords       (N, 2)", WHITE),
    ("patch_classes (N,)", WHITE),
    ("tissue_id    (N,)", WHITE),
    ("slide_width  scalar", LIGHT_GRAY),
    ("slide_height scalar", LIGHT_GRAY),
    ("", WHITE),
    ("252 slides → 252 H5 files", ACCENT_ORANGE),
]
for text, col in h5_lines:
    p2 = tf.add_paragraph()
    p2.text = text
    p2.font.size = Pt(11)
    p2.font.color.rgb = col
    p2.font.name = "Consolas"
    p2.alignment = PP_ALIGN.LEFT

# Bottom: key parameters
add_multiline_text(slide, Inches(0.5), Inches(4.7), Inches(12), Inches(2.5), [
    "40x magnification  |  1024×1024 px patches  |  ≥60% tissue threshold  |  GrandQC segmentation",
    "Patch classifier: MLP (2560→256→128→3)  —  F1: PanNET 97.2%, Normal 96.2%, Stroma 90.9%",
    "3 = RGB color channels (Red, Green, Blue per pixel)",
], font_size=13, color=LIGHT_GRAY, bullet=True)


# ============================================================================
# SLIDE 5: Stage 2 — AutoEncoder
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(8), Inches(0.6),
             "Stage 2: AutoEncoder — Feature Compression", font_size=28, bold=True, color=BLACK)
add_text_box(slide, Inches(0.5), Inches(0.9), Inches(10), Inches(0.4),
             "Why? 1280-d features → overfitting on small dataset (73 patients). Compress to 256-d.",
             font_size=14, color=LIGHT_GRAY)

# Architecture diagram
layers = [
    ("1280", "Input\n(VirChow2)", ACCENT_BLUE, Inches(1.5)),
    ("768", "", MID_GRAY, Inches(1.2)),
    ("512", "", MID_GRAY, Inches(0.9)),
    ("256", "Latent\n(keep this)", ACCENT_GREEN, Inches(0.7)),
    ("512", "", MID_GRAY, Inches(0.9)),
    ("768", "", MID_GRAY, Inches(1.2)),
    ("1280", "Reconstruction", ACCENT_BLUE, Inches(1.5)),
]

x = Inches(0.8)
center_y = Inches(3.8)
gap = Inches(0.15)

for i, (dim, label, col, h) in enumerate(layers):
    bar_top = center_y - h / 2
    bar_w = Inches(1.3)
    card = add_card(slide, x, bar_top, bar_w, h, fill_color=col if i == 3 else DARK_CARD, border_color=BORDER_GRAY)
    tf = card.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = dim
    p.font.size = Pt(16)
    p.font.color.rgb = BLACK
    p.font.bold = True

    if label:
        add_text_box(slide, x, bar_top - Inches(0.5), bar_w, Inches(0.5),
                     label, font_size=10, color=BLACK, alignment=PP_ALIGN.CENTER)

    if i < 6:
        arrow_x = x + bar_w
        add_arrow(slide, arrow_x, center_y - Inches(0.1), Inches(0.15), Inches(0.2),
                  ACCENT_GREEN if i == 2 else MID_GRAY)

    x += bar_w + gap

# Encoder / Decoder labels
add_text_box(slide, Inches(0.8), Inches(5.5), Inches(5), Inches(0.4),
             "◄─── ENCODER ───►", font_size=14, bold=True, color=ACCENT_PURPLE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(6.5), Inches(5.5), Inches(5), Inches(0.4),
             "◄─── DECODER ───►", font_size=14, bold=True, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Training details on the right
details_x = Inches(9.5)
add_text_box(slide, details_x, Inches(1.5), Inches(3.5), Inches(0.4),
             "Training Details", font_size=18, bold=True, color=BLACK)

add_multiline_text(slide, details_x, Inches(2.1), Inches(3.5), Inches(4), [
    "Loss: MSE + variance reg",
    "Optimizer: AdamW",
    "LR: 1e-3, WD: 1e-4",
    "Batch size: 4096",
    "Epochs: 50",
    "Early stopping: patience 10",
    "Split: 90/10 by H5 file",
    "",
    "Each layer:",
    "  Linear → RMSNorm → GELU",
    "  + Dropout(0.2) in encoder",
    "",
    "After training: FREEZE encoder",
    "Use as fixed projector in GNN",
], font_size=12, color=LIGHT_GRAY, bullet=False)


# ============================================================================
# SLIDE 6: Stage 3 — Graph Construction (the key innovation)
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
             "Stage 3: Bipartite Graph Construction", font_size=28, bold=True, color=BLACK)
add_text_box(slide, Inches(0.5), Inches(0.9), Inches(12), Inches(0.4),
             "The core innovation: restrict edges to the tumor–NNP interface",
             font_size=14, color=LIGHT_GRAY)

# Step-by-step pipeline
steps = [
    ("1. Hole Filling", "Fill gaps inside\ntumor regions\n(morphological)", ACCENT_PURPLE),
    ("2. Border Detection", "Find PanNET patches\nwith non-PanNET\n8-neighbors", ACCENT_ORANGE),
    ("3. NNP Inclusion", "Include non-tumor\npatches within\nr hops of border", ACCENT_BLUE),
    ("4. Bipartite Edges", "Connect ONLY\nPanNET ↔ NNP\n(no same-class)", ACCENT_GREEN),
    ("5. Filter Components", "Remove clusters\n< threshold nodes\n(noise removal)", MID_GRAY),
]

card_w = Inches(2.2)
card_h = Inches(1.8)
start_x = Inches(0.4)

for i, (title, desc, col) in enumerate(steps):
    x = start_x + i * (card_w + Inches(0.25))
    card = add_card(slide, x, Inches(1.6), card_w, card_h, border_color=BORDER_GRAY)
    tf = card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(13)
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(11)
    p2.font.color.rgb = LIGHT_GRAY
    p2.alignment = PP_ALIGN.CENTER

    if i < 4:
        add_arrow(slide, x + card_w, Inches(2.2), Inches(0.25), Inches(0.25), MID_GRAY)

# Full Graph vs Bipartite comparison
comp_y = Inches(3.8)
add_text_box(slide, Inches(0.5), comp_y, Inches(5), Inches(0.4),
             "Full Graph (all edges)", font_size=16, bold=True, color=BLACK)
add_text_box(slide, Inches(7), comp_y, Inches(5), Inches(0.4),
             "Bipartite Graph (thesis approach)", font_size=16, bold=True, color=BLACK)

# Full graph visual — all nodes connected
fg_x = Inches(1.0)
fg_y = Inches(4.4)
# PanNET nodes (top row)
for i in range(3):
    add_circle(slide, fg_x + i * Inches(0.8), fg_y, Inches(0.4), ACCENT_GREEN, "T", 9)
# Stroma nodes (bottom row)
for i in range(3):
    add_circle(slide, fg_x + i * Inches(0.8), fg_y + Inches(0.8), Inches(0.4), ACCENT_RED, "S", 9)

add_text_box(slide, fg_x - Inches(0.2), fg_y + Inches(1.4), Inches(3.5), Inches(0.8),
             "All edges: T↔T, S↔S, T↔S\nNoisy — wastes capacity on\nhomogeneous interactions",
             font_size=11, color=ACCENT_RED)

# Bipartite graph visual
bg_x = Inches(7.5)
for i in range(3):
    add_circle(slide, bg_x + i * Inches(0.8), fg_y, Inches(0.4), ACCENT_GREEN, "T", 9)
for i in range(3):
    add_circle(slide, bg_x + i * Inches(0.8), fg_y + Inches(0.8), Inches(0.4), ACCENT_RED, "S", 9)

add_text_box(slide, bg_x - Inches(0.2), fg_y + Inches(1.4), Inches(3.5), Inches(0.8),
             "Only T↔S edges\nForces learning at the\ntumor–NNP interface",
             font_size=11, color=ACCENT_GREEN)

# Bottom insight
card = add_card(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.7), border_color=BORDER_GRAY)
tf = card.text_frame
p = tf.paragraphs[0]
p.text = "Thesis finding: Bipartite consistently outperforms Full Graph across all neighborhood radii (Table 5.5)"
p.font.size = Pt(13)
p.font.color.rgb = ACCENT_GREEN
p.font.bold = True
p.alignment = PP_ALIGN.CENTER


# ============================================================================
# SLIDE 7: Hop Distance / Neighborhood Radius
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
             "Graph Construction: Hop Distance (Radius r)", font_size=28, bold=True, color=BLACK)

# r=1 grid
grids = [
    ("r = 1", "8 neighbors", Inches(0.5)),
    ("r = 2", "24 neighbors", Inches(4.7)),
    ("r = 3", "48 neighbors", Inches(8.9)),
]

for title, subtitle, gx in grids:
    add_text_box(slide, gx, Inches(1.2), Inches(3.5), Inches(0.4),
                 title, font_size=18, bold=True, color=BLACK, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, gx, Inches(1.6), Inches(3.5), Inches(0.3),
                 subtitle, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    grid_size = 7 if "3" in title else (5 if "2" in title else 3)
    cell_size = Inches(0.38) if grid_size <= 5 else Inches(0.3)
    grid_start_x = gx + Inches(1.75) - (grid_size * cell_size) / 2
    grid_start_y = Inches(2.1)

    center = grid_size // 2
    radius = 1 if "1" in title else (2 if "2" in title else 3)

    for row in range(grid_size):
        for col in range(grid_size):
            px = grid_start_x + col * cell_size
            py = grid_start_y + row * cell_size
            dist = max(abs(row - center), abs(col - center))
            if row == center and col == center:
                c = ACCENT_ORANGE  # Center patch
            elif dist <= radius:
                c = ACCENT_GREEN   # Connected neighbor
            else:
                c = DARK_CARD
            add_card(slide, px, py, cell_size - Inches(0.03), cell_size - Inches(0.03), fill_color=c)

# Legend
add_circle(slide, Inches(3), Inches(5.5), Inches(0.25), ACCENT_ORANGE)
add_text_box(slide, Inches(3.3), Inches(5.5), Inches(2), Inches(0.3), "Center patch", font_size=11, color=ACCENT_ORANGE)
add_circle(slide, Inches(5), Inches(5.5), Inches(0.25), ACCENT_GREEN)
add_text_box(slide, Inches(5.3), Inches(5.5), Inches(2), Inches(0.3), "Connected neighbor", font_size=11, color=ACCENT_GREEN)
add_circle(slide, Inches(7.5), Inches(5.5), Inches(0.25), DARK_CARD)
add_text_box(slide, Inches(7.8), Inches(5.5), Inches(2), Inches(0.3), "Out of range", font_size=11, color=MID_GRAY)

add_text_box(slide, Inches(0.5), Inches(6.0), Inches(12), Inches(0.4),
             "Distance metric: Chebyshev (L∞) = max(|Δx|, |Δy|) / patch_size", font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

card = add_card(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.7), border_color=BORDER_GRAY)
tf = card.text_frame
p = tf.paragraphs[0]
p.text = "Thesis finding: Larger radius r consistently improves performance. r=3 is optimal (Table 5.5)"
p.font.size = Pt(13)
p.font.color.rgb = ACCENT_GREEN
p.font.bold = True
p.alignment = PP_ALIGN.CENTER


# ============================================================================
# SLIDE 8: Stage 5 — GIN Architecture
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
             "Stage 5: GIN Model Architecture", font_size=28, bold=True, color=BLACK)

# Architecture flow
arch_steps = [
    ("Frozen\nProjector", "AE Encoder\n1280 → 256", ACCENT_PURPLE, Inches(1.8)),
    ("RMSNorm\n(256)", "Normalize\nlatent features", MID_GRAY, Inches(1.5)),
    ("GINConv\n+ Residual", "MLP: 256→256→256\nSum aggregation", ACCENT_RED, Inches(2.2)),
    ("PanNET-Only\nPooling", "global_add_pool\nover tumor nodes", ACCENT_GREEN, Inches(1.8)),
    ("Regression\nHead", "256→128→1\nHuber loss (δ=2)", ACCENT_ORANGE, Inches(1.8)),
]

y = Inches(1.4)
x = Inches(0.3)
for i, (title, desc, col, w) in enumerate(arch_steps):
    card = add_card(slide, x, y, w, Inches(2.0), border_color=BORDER_GRAY)
    tf = card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(15)
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "\n" + desc
    p2.font.size = Pt(11)
    p2.font.color.rgb = LIGHT_GRAY
    p2.alignment = PP_ALIGN.CENTER

    if i < 4:
        add_arrow(slide, x + w, y + Inches(0.7), Inches(0.3), Inches(0.25), MID_GRAY)
        x += w + Inches(0.3)
    else:
        x += w

# GINConv detail
add_text_box(slide, Inches(0.5), Inches(3.7), Inches(6), Inches(0.4),
             "GINConv Layer Detail", font_size=18, bold=True, color=BLACK)

gin_detail = [
    "h_v^(k) = MLP( (1+ε) · h_v^(k-1) + Σ h_u^(k-1) )",
    "",
    "MLP: Linear(256,256) → RMSNorm → ReLU → Dropout(0.4)",
    "      → Linear(256,256) → RMSNorm → ReLU",
    "",
    "Residual: x = GINConv(dropout(x)) + x",
    "",
    "Why SUM (not mean/max)?",
    "  Counts tumor-NNP contact points — more invasion = more interfaces",
]
add_multiline_text(slide, Inches(0.5), Inches(4.2), Inches(6), Inches(3), gin_detail,
                   font_size=12, color=LIGHT_GRAY)

# Right side: key numbers
add_text_box(slide, Inches(7.5), Inches(3.7), Inches(5), Inches(0.4),
             "Model Parameters", font_size=18, bold=True, color=BLACK)

params = [
    ("AutoEncoder (frozen)", "3,021,312", ACCENT_PURPLE),
    ("GIN (1 layer)", "132,096", ACCENT_RED),
    ("Regression head", "33,665", ACCENT_ORANGE),
    ("Total trainable", "165,761", WHITE),
]

for i, (name, count, col) in enumerate(params):
    y = Inches(4.3) + i * Inches(0.5)
    add_text_box(slide, Inches(7.5), y, Inches(3), Inches(0.4), name, font_size=13, color=col)
    add_text_box(slide, Inches(10.5), y, Inches(2), Inches(0.4), count, font_size=13, color=col, bold=True)

# Training config
add_text_box(slide, Inches(7.5), Inches(6.0), Inches(5), Inches(0.4),
             "Training Config", font_size=14, bold=True, color=BLACK)
add_multiline_text(slide, Inches(7.5), Inches(6.4), Inches(5), Inches(1), [
    "AdamW | LR: 1e-4 | WD: 1e-3 | Batch: 8",
    "Max 200 epochs | Early stop: patience 15",
    "Scheduler: ReduceLROnPlateau (0.8, pat=5)",
], font_size=11, color=LIGHT_GRAY)


# ============================================================================
# SLIDE 9: Why Bipartite + PanNET-only pooling works
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
             "Key Insight: Why Bipartite + PanNET-Only Pooling", font_size=28, bold=True, color=BLACK)

# Message passing diagram
add_text_box(slide, Inches(0.5), Inches(1.2), Inches(5.5), Inches(0.4),
             "Bipartite Message Passing (1 layer)", font_size=18, bold=True, color=BLACK)

# Before
add_text_box(slide, Inches(0.8), Inches(1.8), Inches(2.5), Inches(0.3),
             "Before GNN:", font_size=14, bold=True, color=LIGHT_GRAY)

# Tumor nodes
for i in range(3):
    add_circle(slide, Inches(1.0) + i * Inches(0.9), Inches(2.3), Inches(0.5), ACCENT_GREEN, "T", 11)
# NNP nodes
for i in range(3):
    add_circle(slide, Inches(1.0) + i * Inches(0.9), Inches(3.2), Inches(0.5), ACCENT_RED, "S", 11)

add_text_box(slide, Inches(0.5), Inches(3.9), Inches(3.5), Inches(0.5),
             "T nodes: tumor features only\nS nodes: stroma features only",
             font_size=10, color=LIGHT_GRAY)

# Arrow
add_arrow(slide, Inches(3.8), Inches(2.8), Inches(1.0), Inches(0.4), ACCENT_ORANGE)
add_text_box(slide, Inches(3.8), Inches(2.3), Inches(1.0), Inches(0.4),
             "GIN\nlayer", font_size=11, bold=True, color=ACCENT_ORANGE, alignment=PP_ALIGN.CENTER)

# After
add_text_box(slide, Inches(5.2), Inches(1.8), Inches(3), Inches(0.3),
             "After GNN:", font_size=14, bold=True, color=LIGHT_GRAY)

for i in range(3):
    # Tumor nodes now have mixed colors (absorbed NNP info)
    c = add_circle(slide, Inches(5.5) + i * Inches(0.9), Inches(2.3), Inches(0.5), ACCENT_GREEN, "T+S", 9)
for i in range(3):
    c = add_circle(slide, Inches(5.5) + i * Inches(0.9), Inches(3.2), Inches(0.5), ACCENT_RED, "S+T", 9)

add_text_box(slide, Inches(5.0), Inches(3.9), Inches(3.5), Inches(0.5),
             "T nodes: tumor + NNP context\nS nodes: stroma + tumor context",
             font_size=10, color=LIGHT_GRAY)

# Arrow to pooling
add_arrow(slide, Inches(8.2), Inches(2.8), Inches(0.8), Inches(0.4), ACCENT_ORANGE)

# Pooling box
pool_card = add_card(slide, Inches(9.2), Inches(1.8), Inches(3.5), Inches(2.5), border_color=BORDER_GRAY)
tf = pool_card.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Pool PanNET Only"
p.font.size = Pt(16)
p.font.color.rgb = ACCENT_GREEN
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

pool_lines = [
    "",
    "global_add_pool(T nodes)",
    "",
    "NNP info already absorbed",
    "into T embeddings via",
    "message passing.",
    "",
    "Pooling S nodes would",
    "DILUTE the tumor signal.",
]
for line in pool_lines:
    p2 = tf.add_paragraph()
    p2.text = line
    p2.font.size = Pt(11)
    p2.font.color.rgb = LIGHT_GRAY if "NNP" not in line and "DILUTE" not in line else ACCENT_ORANGE
    p2.alignment = PP_ALIGN.CENTER

# Bottom: comparison
comp_y = Inches(5.0)
comparisons = [
    ("MIL (no graph)", "Treats patches independently\nLoses spatial relationships", ACCENT_RED, "~63% F1"),
    ("Full Graph GNN", "All edges (T↔T, S↔S, T↔S)\nWastes capacity on noise", ACCENT_ORANGE, "~67% F1"),
    ("Bipartite GNN", "Only T↔S edges\nFocused on interface", ACCENT_GREEN, "70.7% F1"),
]

for i, (title, desc, col, f1) in enumerate(comparisons):
    x = Inches(0.5) + i * Inches(4.3)
    card = add_card(slide, x, comp_y, Inches(3.8), Inches(2.0), border_color=BORDER_GRAY)
    tf = card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(11)
    p2.font.color.rgb = LIGHT_GRAY
    p2.alignment = PP_ALIGN.CENTER
    p3 = tf.add_paragraph()
    p3.text = f1
    p3.font.size = Pt(18)
    p3.font.color.rgb = col
    p3.font.bold = True
    p3.alignment = PP_ALIGN.CENTER


# ============================================================================
# SLIDE 10: Results — Thesis Performance
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
             "Results: Thesis Performance Summary", font_size=28, bold=True, color=BLACK)

# Main results table
headers = ["Model", "Loss", "IPS-A F1", "IPS-B F1", "IPS-C F1", "Macro F1", "QWK"]
rows = [
    ("DeepSets", "REG", "47.3", "65.0", "77.1", "63.1", "54.1", MID_GRAY),
    ("CLAM", "CLS", "35.3", "41.0", "67.8", "48.0", "37.0", MID_GRAY),
    ("AB-MIL", "REG", "48.2", "63.4", "78.1", "63.2", "55.5", MID_GRAY),
    ("Patch-GCN (1L)", "REG", "41.1", "65.1", "77.9", "61.4", "50.7", MID_GRAY),
    ("Context-Aware MIL", "REG", "47.1", "62.7", "80.1", "63.3", "62.9", ACCENT_ORANGE),
    ("Bipartite GIN (ours)", "REG", "61.0", "67.0", "84.2", "70.7", "69.1", ACCENT_GREEN),
]

table_top = Inches(1.2)
col_widths = [Inches(2.5), Inches(0.8), Inches(1.2), Inches(1.2), Inches(1.2), Inches(1.5), Inches(1.2)]
row_h = Inches(0.45)

# Header row
x = Inches(1.0)
for j, (header, w) in enumerate(zip(headers, col_widths)):
    add_text_box(slide, x, table_top, w, row_h, header, font_size=12, bold=True, color=ACCENT_ORANGE, alignment=PP_ALIGN.CENTER)
    x += w

# Data rows
for i, (model, loss, ips1, ips2, ips3, macro, qwk, col) in enumerate(rows):
    y = table_top + (i + 1) * row_h
    x = Inches(1.0)
    values = [model, loss, ips1, ips2, ips3, macro, qwk]
    for j, (val, w) in enumerate(zip(values, col_widths)):
        is_best = (i == 5)
        fc = ACCENT_GREEN if is_best else col
        add_text_box(slide, x, y, w, row_h, val, font_size=12, color=fc,
                     bold=(j == 0 or is_best), alignment=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)
        x += w

# Key ablation findings
findings_y = Inches(4.5)
add_text_box(slide, Inches(0.5), findings_y, Inches(12), Inches(0.5),
             "Key Ablation Findings", font_size=20, bold=True, color=BLACK)

findings = [
    ("Bipartite > Full Graph", "T↔S only edges outperform all-edges graphs at every radius", ACCENT_GREEN),
    ("Regression > Classification", "Huber loss >> cross-entropy for ordinal grading", ACCENT_GREEN),
    ("1 layer > 2-3 layers", "Deeper GNNs over-smooth at larger radii", ACCENT_ORANGE),
    ("r=3 > r=2 > r=1", "More spatial context consistently helps", ACCENT_GREEN),
    ("GIN ≈ GATv2", "Sum aggregation matches attention; simpler is better", MID_GRAY),
]

for i, (title, desc, col) in enumerate(findings):
    y = findings_y + Inches(0.5) + i * Inches(0.5)
    add_text_box(slide, Inches(0.8), y, Inches(3.5), Inches(0.4), title, font_size=13, bold=True, color=col)
    add_text_box(slide, Inches(4.5), y, Inches(8), Inches(0.4), desc, font_size=12, color=LIGHT_GRAY)


# ============================================================================
# SLIDE 11: Dataset & Evaluation Protocol
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
             "Dataset & Evaluation Protocol", font_size=28, bold=True, color=BLACK)

# Dataset stats
add_text_box(slide, Inches(0.5), Inches(1.2), Inches(5), Inches(0.4),
             "Dataset: Koç University Hospital", font_size=18, bold=True, color=BLACK)

stats = [
    ("73 patients", "252 H&E-stained WSIs"),
    ("61 patients with 3 WSIs", "→ used for IPS evaluation"),
    ("12 patients with <3 WSIs", "→ training only (never evaluated)"),
]
for i, (stat, desc) in enumerate(stats):
    y = Inches(1.8) + i * Inches(0.5)
    add_text_box(slide, Inches(0.8), y, Inches(2.5), Inches(0.4), stat, font_size=14, bold=True, color=WHITE)
    add_text_box(slide, Inches(3.3), y, Inches(3), Inches(0.4), desc, font_size=12, color=LIGHT_GRAY)

# Grade distribution
add_text_box(slide, Inches(0.5), Inches(3.5), Inches(5), Inches(0.4),
             "Grade Distribution (WSI-level)", font_size=16, bold=True, color=BLACK)

grade_dist = [("G1", "64", ACCENT_GREEN), ("G2", "43", ACCENT_GREEN),
              ("G3", "36", ACCENT_ORANGE), ("G4", "34", ACCENT_ORANGE), ("G5", "20", ACCENT_RED)]
for i, (g, count, col) in enumerate(grade_dist):
    x = Inches(0.8) + i * Inches(1.0)
    card = add_card(slide, x, Inches(4.0), Inches(0.8), Inches(0.8), border_color=BORDER_GRAY)
    tf = card.text_frame
    p = tf.paragraphs[0]
    p.text = g
    p.font.size = Pt(14)
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = count
    p2.font.size = Pt(18)
    p2.font.color.rgb = WHITE
    p2.font.bold = True
    p2.alignment = PP_ALIGN.CENTER

# IPS distribution
add_text_box(slide, Inches(0.5), Inches(5.2), Inches(5), Inches(0.4),
             "IPS Distribution (patient-level)", font_size=16, bold=True, color=BLACK)

ips_dist = [("IPS-A", "12", ACCENT_GREEN), ("IPS-B", "21", ACCENT_ORANGE), ("IPS-C", "28", ACCENT_RED)]
for i, (name, count, col) in enumerate(ips_dist):
    x = Inches(0.8) + i * Inches(1.5)
    card = add_card(slide, x, Inches(5.7), Inches(1.2), Inches(0.8), border_color=BORDER_GRAY)
    tf = card.text_frame
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(14)
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = f"n={count}"
    p2.font.size = Pt(16)
    p2.font.color.rgb = WHITE
    p2.font.bold = True
    p2.alignment = PP_ALIGN.CENTER

# CV protocol on the right
add_text_box(slide, Inches(7), Inches(1.2), Inches(5.5), Inches(0.4),
             "4-Fold Cross-Validation", font_size=18, bold=True, color=BLACK)

cv_lines = [
    "s1, s2, s3, s4: rotate as test/val",
    "s5: always training (incomplete patients)",
    "",
    "Per fold: test=s[i], val=s[(i+1)%4]",
    "          train=remaining + s5",
    "",
    "25 random seeds per configuration",
    "→ robust statistical comparison",
    "→ Wilcoxon signed-rank test",
    "",
    "Metrics:",
    "  • Per-class F1 (IPS-A, B, C)",
    "  • Macro F1 (unweighted mean)",
    "  • QWK (penalizes distant errors)",
    "  • MAE (mean absolute error)",
]
add_multiline_text(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(5), cv_lines,
                   font_size=12, color=LIGHT_GRAY)


# ============================================================================
# SLIDE 12: Patient-Level Aggregation
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
             "Patient-Level IPS Prediction Pipeline", font_size=28, bold=True, color=BLACK)

# Flow: 3 slides → 3 grades → sum → IPS
slides_data = [
    ("Slide 1", "Grade: 4"),
    ("Slide 2", "Grade: 3"),
    ("Slide 3", "Grade: 5"),
]

for i, (name, grade) in enumerate(slides_data):
    y = Inches(1.5) + i * Inches(1.5)

    # WSI box
    card = add_card(slide, Inches(0.5), y, Inches(2.0), Inches(1.0), border_color=BORDER_GRAY)
    tf = card.text_frame
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(14)
    p.font.color.rgb = ACCENT_BLUE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Arrow
    add_arrow(slide, Inches(2.6), y + Inches(0.3), Inches(0.5), Inches(0.3), MID_GRAY)

    # Graph box
    card = add_card(slide, Inches(3.3), y, Inches(1.8), Inches(1.0), border_color=BORDER_GRAY)
    tf = card.text_frame
    p = tf.paragraphs[0]
    p.text = "Bipartite\nGraph"
    p.font.size = Pt(12)
    p.font.color.rgb = ACCENT_GREEN
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Arrow
    add_arrow(slide, Inches(5.3), y + Inches(0.3), Inches(0.5), Inches(0.3), MID_GRAY)

    # GIN box
    card = add_card(slide, Inches(6.0), y, Inches(1.5), Inches(1.0), border_color=BORDER_GRAY)
    tf = card.text_frame
    p = tf.paragraphs[0]
    p.text = "GIN\nModel"
    p.font.size = Pt(12)
    p.font.color.rgb = ACCENT_RED
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Arrow
    add_arrow(slide, Inches(7.7), y + Inches(0.3), Inches(0.5), Inches(0.3), MID_GRAY)

    # Grade prediction
    card = add_card(slide, Inches(8.4), y, Inches(1.5), Inches(1.0), border_color=BORDER_GRAY)
    tf = card.text_frame
    p = tf.paragraphs[0]
    p.text = grade
    p.font.size = Pt(16)
    p.font.color.rgb = ACCENT_ORANGE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

# Sum arrow
add_text_box(slide, Inches(10.2), Inches(1.5), Inches(0.5), Inches(4.0),
             "}\n\nΣ", font_size=36, color=WHITE, alignment=PP_ALIGN.CENTER)

# IPS result
card = add_card(slide, Inches(10.8), Inches(2.3), Inches(2.2), Inches(2.0), border_color=BORDER_GRAY)
tf = card.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Sum = 12"
p.font.size = Pt(18)
p.font.color.rgb = WHITE
p.font.bold = True
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "\nIPS-C"
p2.font.size = Pt(24)
p2.font.color.rgb = ACCENT_RED
p2.font.bold = True
p2.alignment = PP_ALIGN.CENTER
p3 = tf.add_paragraph()
p3.text = "Highly\nInfiltrative"
p3.font.size = Pt(12)
p3.font.color.rgb = LIGHT_GRAY
p3.alignment = PP_ALIGN.CENTER

# Bottom explanation
add_multiline_text(slide, Inches(0.5), Inches(5.5), Inches(12), Inches(1.5), [
    "Each WSI → graph → GIN → predicted grade (1-5 continuous → rounded to integer)",
    "Sum 3 slide grades: [3-6] = IPS-A  |  [7-9] = IPS-B  |  [10-15] = IPS-C",
    "Only patients with exactly 3 slides are evaluated (61 patients)",
], font_size=13, color=LIGHT_GRAY, bullet=True)


# ============================================================================
# Save
# ============================================================================
output_path = "/Users/sinan/gnn/pannet_pipeline/PanNET_Pipeline_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
