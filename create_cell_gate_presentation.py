"""
Generate a compact PowerPoint presentation explaining the Cell-Conditioned
Gate mechanism on top of the base bipartite GIN model.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors — minimal, no unnecessary coloring
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
DARK_GRAY = RGBColor(0x44, 0x44, 0x44)
MID_GRAY = RGBColor(0x77, 0x77, 0x77)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
TABLE_HEADER_BG = RGBColor(0x33, 0x33, 0x33)
TABLE_ALT_BG = RGBColor(0xF7, 0xF7, 0xF7)

# Accent — used sparingly
TUMOR_COLOR = RGBColor(0xC0, 0x39, 0x2B)      # Red for tumor/PanNET
NNP_COLOR = RGBColor(0x29, 0x80, 0xB9)        # Blue for NNP
GATE_COLOR = RGBColor(0xE6, 0x7E, 0x22)       # Orange for gate mechanism

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_text(slide, left, top, width, height, text, size=18, color=BLACK,
             bold=False, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return box


def add_rich_text(slide, left, top, width, height, lines, font="Calibri"):
    """lines: list of (text, size, color, bold)"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, size, color, bold) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font
        p.space_after = Pt(4)
    return box


def add_box(slide, left, top, width, height, fill=LIGHT_GRAY, border=BORDER_GRAY):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    ln = shape.line
    ln.color.rgb = border
    ln.width = Pt(1)
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=MID_GRAY, width=Pt(2)):
    connector = slide.shapes.add_connector(
        1, x1, y1, x2, y2  # 1 = straight connector
    )
    connector.line.color.rgb = color
    connector.line.width = width
    # Add arrowhead
    connector.end_x = x2
    connector.end_y = y2
    return connector


def add_table(slide, left, top, width, rows_data, col_widths, header=True):
    """Add a simple table. rows_data: list of lists of strings."""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, Inches(0.4 * n_rows))
    table = table_shape.table

    for i, cw in enumerate(col_widths):
        table.columns[i].width = cw

    for r, row in enumerate(rows_data):
        for c, cell_text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = cell_text
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.name = "Calibri"
                paragraph.alignment = PP_ALIGN.LEFT

                if r == 0 and header:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                else:
                    paragraph.font.color.rgb = BLACK

            # Styling
            if r == 0 and header:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER_BG
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT_BG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE

    return table_shape


# =============================================================================
# SLIDE 1: Title
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)

add_text(slide, Inches(1), Inches(2.5), Inches(11), Inches(1),
         "Cell-Conditioned Gate Mechanism", size=40, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.8),
         "Improving Bipartite GIN with Cell Composition Information", size=22,
         color=MID_GRAY, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(4.8), Inches(11), Inches(0.6),
         "PanNET Infiltration Pattern Scoring", size=16,
         color=MID_GRAY, align=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 2: What We Have — Cell Types from HoVer-Net
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.6),
         "What We Have: Cell Composition Per Patch", size=28, bold=True)
add_text(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.5),
         "HoVer-Net (PanNuke) detected ~57M nuclei across 184 slides, classified into 6 types, collapsed to 4 channels:",
         size=14, color=DARK_GRAY)

# Cell types table
cell_type_data = [
    ["Channel", "Contents", "Source PanNuke Types", "Total Cells"],
    ["[0] neoplastic", "Tumor cells", "Type 1", "20.5M"],
    ["[1] inflammatory", "Immune cells", "Type 2", "3.9M"],
    ["[2] other", "Connective + dead + macrophages", "Types 3, 4, 5", "32.6M"],
    ["[3] reserved", "Always 0", "—", "—"],
]
add_table(slide, Inches(0.6), Inches(1.6), Inches(8.5), cell_type_data,
          [Inches(1.8), Inches(3.2), Inches(2.0), Inches(1.5)])

# Right side: what this means
add_box(slide, Inches(9.5), Inches(1.6), Inches(3.3), Inches(2.8), fill=LIGHT_GRAY, border=BORDER_GRAY)
add_rich_text(slide, Inches(9.7), Inches(1.7), Inches(2.9), Inches(2.6), [
    ("Each patch gets a 4-dim vector:", 13, BLACK, True),
    ("", 6, BLACK, False),
    ("Example — immune-rich NNP patch:", 12, DARK_GRAY, False),
    ("[  2,  45,  12,  0 ]", 14, BLACK, True),
    ("  2 tumor, 45 immune, 12 other", 11, MID_GRAY, False),
    ("", 6, BLACK, False),
    ("Example — dense tumor patch:", 12, DARK_GRAY, False),
    ("[ 120,  3,   8,  0 ]", 14, BLACK, True),
    ("  120 tumor, 3 immune, 8 other", 11, MID_GRAY, False),
])

# Coverage subtitle
add_text(slide, Inches(0.6), Inches(4.7), Inches(12), Inches(0.5),
         "Coverage: 184 / 252 slides (73%), from 51 / 73 patients. Missing slides get zero-filled vectors.",
         size=14, color=DARK_GRAY)

# Coverage by fold table
fold_data = [
    ["Fold", "Slides With", "Slides Without", "Patients OK", "Patients Missing"],
    ["s1", "43", "17", "15", "6  (#32, #45, #47, #52, #68, #69)"],
    ["s2", "38", "21", "13", "5  (#36, #37, #44, #57, #72)"],
    ["s3", "39", "17", "13", "6  (#16, #25, #41, #48, #49, #70)"],
    ["s4", "44", "13", "15", "5  (#43, #46, #59, #71, #73)"],
    ["s5 (train only)", "20", "0", "12", "0"],
]
add_table(slide, Inches(0.6), Inches(5.3), Inches(10), fold_data,
          [Inches(1.5), Inches(1.3), Inches(1.5), Inches(1.3), Inches(4.4)])


# =============================================================================
# SLIDE 3: Base GIN vs Cell-Conditioned GIN — Architecture Comparison
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.6),
         "What Changes: Base GIN vs Cell-Conditioned GIN", size=28, bold=True)

# --- Base GIN (top row) ---
add_text(slide, Inches(0.6), Inches(1.2), Inches(4), Inches(0.4),
         "Base GIN (cell_info_mode = none)", size=16, bold=True, color=DARK_GRAY)

box_y = Inches(1.7)
box_h = Inches(1.2)
boxes_base = [
    ("RMSNorm\n(256)", Inches(0.6), Inches(1.8)),
    ("GINConv\n+ Residual", Inches(2.7), Inches(2.4)),
    ("PanNET-Only\nPooling", Inches(5.5), Inches(2.0)),
    ("Regression\nHead", Inches(7.9), Inches(1.8)),
]

for label, left, w in boxes_base:
    shape = add_box(slide, left, box_y, w, box_h, fill=LIGHT_GRAY, border=BORDER_GRAY)
    shape.text_frame.word_wrap = True
    for i, line in enumerate(label.split("\n")):
        if i == 0:
            p = shape.text_frame.paragraphs[0]
        else:
            p = shape.text_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.color.rgb = BLACK
        p.font.bold = (i == 0)
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
    shape.text_frame.paragraphs[0].space_before = Pt(12)

# Arrows between base boxes
for i in range(len(boxes_base) - 1):
    _, l1, w1 = boxes_base[i]
    _, l2, w2 = boxes_base[i + 1]
    x1 = l1 + w1
    x2 = l2
    y_mid = box_y + box_h // 2
    add_arrow(slide, x1, y_mid, x2, y_mid, color=MID_GRAY, width=Pt(1.5))

# Description under base GINConv
add_text(slide, Inches(2.5), Inches(2.95), Inches(2.8), Inches(0.7),
         "new_h(v) = MLP( h(v) + SUM[ h(u) ] ) + h(v)\nAll neighbors weighted equally",
         size=10, color=MID_GRAY)

# --- Cell-Conditioned GIN (bottom row) ---
add_text(slide, Inches(0.6), Inches(3.8), Inches(5), Inches(0.4),
         "Cell-Conditioned GIN (cell_info_mode = gate)", size=16, bold=True, color=DARK_GRAY)

box_y2 = Inches(4.3)
boxes_gate = [
    ("RMSNorm\n(256)", Inches(0.6), Inches(1.8)),
    ("CellConditioned\nConv + Residual", Inches(2.7), Inches(2.4)),
    ("PanNET-Only\nPooling", Inches(5.5), Inches(2.0)),
    ("Regression\nHead", Inches(7.9), Inches(1.8)),
]

for idx, (label, left, w) in enumerate(boxes_gate):
    if idx == 1:
        # Highlight the changed box
        shape = add_box(slide, left, box_y2, w, box_h, fill=RGBColor(0xFD, 0xF2, 0xE9), border=GATE_COLOR)
    else:
        shape = add_box(slide, left, box_y2, w, box_h, fill=LIGHT_GRAY, border=BORDER_GRAY)
    shape.text_frame.word_wrap = True
    for i, line in enumerate(label.split("\n")):
        if i == 0:
            p = shape.text_frame.paragraphs[0]
        else:
            p = shape.text_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.color.rgb = BLACK
        p.font.bold = (i == 0)
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
    shape.text_frame.paragraphs[0].space_before = Pt(12)

# Arrows between gate boxes
for i in range(len(boxes_gate) - 1):
    _, l1, w1 = boxes_gate[i]
    _, l2, w2 = boxes_gate[i + 1]
    x1 = l1 + w1
    x2 = l2
    y_mid = box_y2 + box_h // 2
    add_arrow(slide, x1, y_mid, x2, y_mid, color=MID_GRAY, width=Pt(1.5))

# Cell info arrow feeding into the gate conv
cell_info_box = add_box(slide, Inches(3.0), Inches(5.7), Inches(1.9), Inches(0.7),
                        fill=RGBColor(0xFD, 0xF2, 0xE9), border=GATE_COLOR)
cell_info_box.text_frame.paragraphs[0].text = "cell_information"
cell_info_box.text_frame.paragraphs[0].font.size = Pt(11)
cell_info_box.text_frame.paragraphs[0].font.bold = True
cell_info_box.text_frame.paragraphs[0].font.name = "Calibri"
cell_info_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
p2 = cell_info_box.text_frame.add_paragraph()
p2.text = "(N, 4) per node"
p2.font.size = Pt(10)
p2.font.color.rgb = MID_GRAY
p2.font.name = "Calibri"
p2.alignment = PP_ALIGN.CENTER

# Arrow from cell_info up to gate conv
add_arrow(slide, Inches(3.95), Inches(5.7), Inches(3.95), box_y2 + box_h,
          color=GATE_COLOR, width=Pt(1.5))

# "ONLY THIS CHANGES" label
add_text(slide, Inches(5.3), Inches(4.1), Inches(2.5), Inches(0.3),
         "only this box changes", size=11, color=GATE_COLOR, bold=True)

# Comparison table on right
comp_data = [
    ["Component", "Base GIN", "Gate GIN"],
    ["Normalization", "RMSNorm(256)", "Same"],
    ["Message passing", "GINConv: sum all\nneighbors equally", "CellConditionedConv:\nper-edge gating"],
    ["Pooling", "global_add_pool\ntumor only", "Same"],
    ["Regression head", "256 -> 128 -> 1", "Same"],
    ["New parameters", "—", "~200K (gate MLP\n+ msg + update)"],
]
add_table(slide, Inches(8.0), Inches(1.2), Inches(4.8), comp_data,
          [Inches(1.5), Inches(1.5), Inches(1.8)])


# =============================================================================
# SLIDE 4: How the Gate Works — Step by Step
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.6),
         "How the Gate Works: Per-Edge Message Modulation", size=28, bold=True)

# Left side: diagram
# Tumor node (target)
tumor_box = add_box(slide, Inches(0.8), Inches(2.8), Inches(1.6), Inches(1.4),
                    fill=RGBColor(0xFC, 0xE4, 0xE4), border=TUMOR_COLOR)
tumor_box.text_frame.paragraphs[0].text = "Tumor patch (v)"
tumor_box.text_frame.paragraphs[0].font.size = Pt(11)
tumor_box.text_frame.paragraphs[0].font.bold = True
tumor_box.text_frame.paragraphs[0].font.name = "Calibri"
tumor_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
p = tumor_box.text_frame.add_paragraph()
p.text = "h(v) = 256-dim"
p.font.size = Pt(10)
p.font.color.rgb = MID_GRAY
p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER
p = tumor_box.text_frame.add_paragraph()
p.text = "cell(v) = [120, 3, 8, 0]"
p.font.size = Pt(10)
p.font.color.rgb = TUMOR_COLOR
p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER

# NNP node (source)
nnp_box = add_box(slide, Inches(0.8), Inches(1.0), Inches(1.6), Inches(1.4),
                  fill=RGBColor(0xE0, 0xEE, 0xF7), border=NNP_COLOR)
nnp_box.text_frame.paragraphs[0].text = "NNP patch (u)"
nnp_box.text_frame.paragraphs[0].font.size = Pt(11)
nnp_box.text_frame.paragraphs[0].font.bold = True
nnp_box.text_frame.paragraphs[0].font.name = "Calibri"
nnp_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
p = nnp_box.text_frame.add_paragraph()
p.text = "h(u) = 256-dim"
p.font.size = Pt(10)
p.font.color.rgb = MID_GRAY
p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER
p = nnp_box.text_frame.add_paragraph()
p.text = "cell(u) = [2, 45, 12, 0]"
p.font.size = Pt(10)
p.font.color.rgb = NNP_COLOR
p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER

# Arrow from NNP to Tumor
add_arrow(slide, Inches(1.6), Inches(2.4), Inches(1.6), Inches(2.8),
          color=BLACK, width=Pt(2))

# Right side: step-by-step computation
step_x = Inches(3.2)

# Step 1: Concatenate cell info
s1_box = add_box(slide, step_x, Inches(1.0), Inches(4.2), Inches(1.1),
                 fill=LIGHT_GRAY, border=BORDER_GRAY)
add_rich_text(slide, step_x + Inches(0.15), Inches(1.05), Inches(3.9), Inches(1.0), [
    ("Step 1: Concatenate cell compositions", 12, BLACK, True),
    ("cell_pair = [ cell(u) || cell(v) ]", 12, DARK_GRAY, False),
    ("= [ 2, 45, 12, 0,  120, 3, 8, 0 ]      (8-dim vector)", 11, MID_GRAY, False),
])

# Step 2: Compute gate
s2_box = add_box(slide, step_x, Inches(2.3), Inches(4.2), Inches(1.3),
                 fill=RGBColor(0xFD, 0xF2, 0xE9), border=GATE_COLOR)
add_rich_text(slide, step_x + Inches(0.15), Inches(2.35), Inches(3.9), Inches(1.2), [
    ("Step 2: Compute gate (learned)", 12, BLACK, True),
    ("gate = sigmoid( MLP( cell_pair ) )", 12, DARK_GRAY, False),
    ("MLP: Linear(8, 256) -> ReLU -> Linear(256, 256) -> Sigmoid", 10, MID_GRAY, False),
    ("= 256 values, each in [0, 1]  (per-channel on/off)", 11, MID_GRAY, False),
])

# Step 3: Gate the message
s3_box = add_box(slide, step_x, Inches(3.8), Inches(4.2), Inches(1.1),
                 fill=LIGHT_GRAY, border=BORDER_GRAY)
add_rich_text(slide, step_x + Inches(0.15), Inches(3.85), Inches(3.9), Inches(1.0), [
    ("Step 3: Gate the message", 12, BLACK, True),
    ("msg = Linear( h(u) ) * gate", 12, DARK_GRAY, False),
    ("Element-wise: 256-dim message modulated by 256-dim gate", 11, MID_GRAY, False),
])

# Step 4: Aggregate + Update
s4_box = add_box(slide, step_x, Inches(5.1), Inches(4.2), Inches(1.3),
                 fill=LIGHT_GRAY, border=BORDER_GRAY)
add_rich_text(slide, step_x + Inches(0.15), Inches(5.15), Inches(3.9), Inches(1.2), [
    ("Step 4: Aggregate all neighbors + Update", 12, BLACK, True),
    ("agg = SUM[ gated messages from all neighbors ]", 12, DARK_GRAY, False),
    ("new_h(v) = MLP_update( [ h(v) || agg ] ) + h(v)", 12, DARK_GRAY, False),
    ("512-dim -> 256-dim + residual", 11, MID_GRAY, False),
])

# Right side: biological interpretation
interp_box = add_box(slide, Inches(8.0), Inches(1.0), Inches(4.8), Inches(5.4),
                     fill=LIGHT_GRAY, border=BORDER_GRAY)
add_rich_text(slide, Inches(8.2), Inches(1.1), Inches(4.4), Inches(5.2), [
    ("Biological Interpretation", 16, BLACK, True),
    ("", 8, BLACK, False),
    ("The gate learns WHICH cell interactions", 13, DARK_GRAY, False),
    ("matter for infiltration scoring:", 13, DARK_GRAY, False),
    ("", 8, BLACK, False),
    ("High gate (pass message through):", 13, BLACK, True),
    ("Immune-rich NNP next to tumor-dense", 12, DARK_GRAY, False),
    ("patch = active infiltration signal.", 12, DARK_GRAY, False),
    ("The model should pay attention to this.", 12, DARK_GRAY, False),
    ("", 8, BLACK, False),
    ("Low gate (dampen message):", 13, BLACK, True),
    ("Cell-sparse connective patch next to", 12, DARK_GRAY, False),
    ("tumor = less informative boundary.", 12, DARK_GRAY, False),
    ("The model can safely ignore this.", 12, DARK_GRAY, False),
    ("", 8, BLACK, False),
    ("Key difference from base GIN:", 13, BLACK, True),
    ("Base GIN: all neighbors weighted equally.", 12, DARK_GRAY, False),
    ("Gate GIN: neighbors weighted by their", 12, DARK_GRAY, False),
    ("cell composition relative to the target.", 12, DARK_GRAY, False),
    ("", 8, BLACK, False),
    ("In base GIN:", 13, BLACK, True),
    ("  new_h(v) = MLP( h(v) + SUM[h(u)] )", 11, MID_GRAY, False),
    ("", 4, BLACK, False),
    ("In gate GIN:", 13, BLACK, True),
    ("  new_h(v) = MLP( h(v) || SUM[msg(u)*gate(u,v)] )", 11, MID_GRAY, False),
])


# =============================================================================
# SLIDE 5: Handling Missing Data + Expected Outcome
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.6),
         "Handling Missing Data & Experiment Setup", size=28, bold=True)

# Left: missing data handling
add_box(slide, Inches(0.6), Inches(1.3), Inches(5.8), Inches(2.8), fill=LIGHT_GRAY, border=BORDER_GRAY)
add_rich_text(slide, Inches(0.8), Inches(1.4), Inches(5.4), Inches(2.6), [
    ("Missing Cell Data Strategy", 16, BLACK, True),
    ("", 6, BLACK, False),
    ("67 / 251 graphs (27%) have no cell counts.", 13, DARK_GRAY, False),
    ("These get zero-filled: cell(v) = [0, 0, 0, 0]", 13, DARK_GRAY, False),
    ("", 6, BLACK, False),
    ("What happens with zeros in the gate:", 13, BLACK, True),
    ("cell_pair = [0,0,0,0, 0,0,0,0]  (all zeros)", 12, MID_GRAY, False),
    ("gate = sigmoid(MLP([0,...,0]))", 12, MID_GRAY, False),
    ("", 4, BLACK, False),
    ("The gate learns a default pass-through behavior", 13, DARK_GRAY, False),
    ("for zero inputs, effectively falling back to the", 13, DARK_GRAY, False),
    ("base GIN behavior for those patches.", 13, DARK_GRAY, False),
])

# Right: experiment setup
add_box(slide, Inches(6.8), Inches(1.3), Inches(5.8), Inches(2.8), fill=LIGHT_GRAY, border=BORDER_GRAY)
add_rich_text(slide, Inches(7.0), Inches(1.4), Inches(5.4), Inches(2.6), [
    ("Experiment Setup", 16, BLACK, True),
    ("", 6, BLACK, False),
    ("Same as base GIN evaluation:", 13, DARK_GRAY, False),
    ("", 4, BLACK, False),
    ("  4-fold CV (s1-s4 rotate test/val, s5 = train)", 12, DARK_GRAY, False),
    ("  5 seeds per fold = 20 training runs", 12, DARK_GRAY, False),
    ("  1 GIN layer, radius = 3 (best config)", 12, DARK_GRAY, False),
    ("  Huber loss (delta = 2), AdamW, lr = 1e-4", 12, DARK_GRAY, False),
    ("  Early stopping patience = 15", 12, DARK_GRAY, False),
    ("  Patient-level IPS evaluation", 12, DARK_GRAY, False),
    ("", 6, BLACK, False),
    ("Only change: GINConv -> CellConditionedConv", 13, BLACK, True),
    ("Everything else is identical.", 13, DARK_GRAY, False),
])

# Bottom: what we expect
add_box(slide, Inches(0.6), Inches(4.5), Inches(12), Inches(2.5), fill=LIGHT_GRAY, border=BORDER_GRAY)
add_rich_text(slide, Inches(0.8), Inches(4.6), Inches(11.6), Inches(2.3), [
    ("What We Expect", 16, BLACK, True),
    ("", 6, BLACK, False),
    ("Hypothesis: cell composition at the tumor-NNP boundary is informative for infiltration grading.", 14, DARK_GRAY, False),
    ("", 4, BLACK, False),
    ("Higher-grade tumors (4-5) show irregular boundaries with immune cell infiltration into the tumor.", 13, DARK_GRAY, False),
    ("Lower-grade tumors (1-2) have clean borders with few immune cells crossing the boundary.", 13, DARK_GRAY, False),
    ("The gate mechanism lets the model distinguish between these patterns at each edge.", 13, DARK_GRAY, False),
    ("", 6, BLACK, False),
    ("Target to beat: Base bipartite GIN at 70.7% macro F1, 69.1 QWK (Ozates thesis, Table 5.5)", 14, BLACK, True),
])


# =============================================================================
# SLIDE 6: Results — Gate vs Baseline
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.6),
         "Results: Gate vs Baseline (20 runs each, 4 folds x 5 seeds)", size=28, bold=True)

# Main comparison table
results_data = [
    ["Metric", "Base GIN (none)", "Cell Gate", "Difference"],
    ["Macro F1", "64.45 +/- 8.43", "60.32 +/- 9.51", "-4.13"],
    ["Weighted F1", "68.50 +/- 7.65", "64.16 +/- 9.98", "-4.34"],
    ["QWK", "67.86 +/- 8.59", "59.16 +/- 11.45", "-8.70"],
    ["F1 IPS-A", "57.32 +/- 14.45", "51.37 +/- 12.47", "-5.95"],
    ["F1 IPS-B", "54.95 +/- 10.50", "53.98 +/- 12.95", "-0.97"],
    ["F1 IPS-C", "81.09 +/- 7.99", "75.62 +/- 9.95", "-5.47"],
]
add_table(slide, Inches(0.6), Inches(1.3), Inches(9), results_data,
          [Inches(1.8), Inches(2.2), Inches(2.2), Inches(1.5)])

# Per-fold breakdown
add_text(slide, Inches(0.6), Inches(4.2), Inches(12), Inches(0.5),
         "Per-Fold Macro F1 (gate)", size=16, bold=True, color=DARK_GRAY)

fold_data = [
    ["Fold", "Seed 42", "Seed 128", "Seed 777", "Seed 1234", "Seed 5", "Fold Avg"],
    ["Fold 0", "50.00", "80.08", "62.30", "63.85", "84.71", "68.19"],
    ["Fold 1", "51.94", "51.10", "65.75", "55.16", "58.33", "56.46"],
    ["Fold 2", "65.56", "65.56", "65.56", "56.72", "65.56", "63.79"],
    ["Fold 3", "46.65", "57.94", "55.91", "49.84", "53.95", "52.86"],
]
add_table(slide, Inches(0.6), Inches(4.8), Inches(10), fold_data,
          [Inches(1.0), Inches(1.2), Inches(1.2), Inches(1.2), Inches(1.2), Inches(1.0), Inches(1.2)])

# Analysis box on right
add_box(slide, Inches(9.3), Inches(1.3), Inches(3.5), Inches(3.2), fill=LIGHT_GRAY, border=BORDER_GRAY)
add_rich_text(slide, Inches(9.5), Inches(1.4), Inches(3.1), Inches(3.0), [
    ("Analysis", 16, BLACK, True),
    ("", 6, BLACK, False),
    ("Gate performs worse than baseline", 13, TUMOR_COLOR, True),
    ("across all metrics.", 13, TUMOR_COLOR, False),
    ("", 6, BLACK, False),
    ("Possible reasons:", 13, BLACK, True),
    ("", 4, BLACK, False),
    ("1. 27% missing cell data (zeros)", 12, DARK_GRAY, False),
    ("   adds noise to gate learning", 12, MID_GRAY, False),
    ("", 4, BLACK, False),
    ("2. Only 20 runs (5 seeds) — high", 12, DARK_GRAY, False),
    ("   variance, need 25 seeds", 12, MID_GRAY, False),
    ("", 4, BLACK, False),
    ("3. ~200K extra parameters may", 12, DARK_GRAY, False),
    ("   overfit on small dataset", 12, MID_GRAY, False),
    ("", 4, BLACK, False),
    ("4. Cell counts may not add signal", 12, DARK_GRAY, False),
    ("   beyond what VirChow2 already", 12, MID_GRAY, False),
    ("   captures in the 256-d features", 12, MID_GRAY, False),
])


# Save
output_path = "Cell_Gate_Mechanism.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
